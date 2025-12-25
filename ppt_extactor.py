import streamlit as st
import os
import json
import io
import time
import shutil
from PIL import Image

# --- LIBRARIES FOR EXTRACTION ---
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import pytesseract
    import requests
except ImportError:
    st.error("Missing Extraction Libraries. Run: pip install python-pptx pytesseract Pillow requests")

# --- LIBRARIES FOR RAG (VECTOR DB) ---
try:
    import chromadb
    from chromadb.utils import embedding_functions
    from sentence_transformers import SentenceTransformer
except ImportError:
    st.error("Missing RAG Libraries. Run: pip install chromadb sentence-transformers")

# --- LIBRARY FOR PPT RENDERING (WINDOWS ONLY) ---
try:
    import win32com.client
    import pythoncom
except ImportError:
    # We don't error immediately, but rendering features won't work without pywin32
    pass

# ==========================================
# 1. CONFIGURATION & SETUP
# ==========================================

st.set_page_config(page_title="PPT RAG Explorer", layout="wide")

# Default Constants
DEFAULT_TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
DEFAULT_API_KEY = "sk-f7d101c4a97246318ab270f5d67abfdd"
DEEPSEEK_URL = "https://api.deepseek.com/chat/completions"
EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2" # Small, fast, effective for local use

# Initialize Session State
if "extracted_data" not in st.session_state:
    st.session_state.extracted_data = None
if "vector_db_ready" not in st.session_state:
    st.session_state.vector_db_ready = False
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "active_slide_id" not in st.session_state:
    st.session_state.active_slide_id = None
if "current_ppt_path" not in st.session_state:
    st.session_state.current_ppt_path = None

# ==========================================
# 2. HELPER: SLIDE RENDERING (VISUAL)
# ==========================================

def render_slide_as_image(ppt_path, slide_index):
    """
    Uses Microsoft PowerPoint (via COM) to export a specific slide as an image.
    ppt_path: Absolute path to the PPTX file.
    slide_index: 1-based index of the slide.
    Returns: Path to the generated image or None if failed.
    """
    try:
        # Initialize COM library (needed for Streamlit threads)
        pythoncom.CoInitialize()
        
        # Connect to PowerPoint
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        # ppt_app.Visible = 1  # Uncomment if debugging, otherwise keep hidden
        
        # Open Presentation (Read-only, No Window)
        presentation = ppt_app.Presentations.Open(ppt_path, WithWindow=False)
        
        # Output File Path
        output_img_path = os.path.abspath(f"temp_slide_{slide_index}.jpg")
        
        # Export Slide (Index is 1-based in COM)
        presentation.Slides(slide_index).Export(output_img_path, "JPG")
        
        presentation.Close()
        return output_img_path
        
    except Exception as e:
        print(f"Rendering Error: {e}")
        return None
    finally:
        # Ensure COM is uninitialized
        pythoncom.CoUninitialize()

# ==========================================
# 3. EXTRACTION LOGIC
# ==========================================

def get_shape_text(shape):
    text_parts = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            p_text = "".join(run.text for run in paragraph.runs)
            if p_text.strip():
                text_parts.append(p_text)
    if shape.has_table:
        for row in shape.table.rows:
            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_cells:
                text_parts.append(" | ".join(row_cells))
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            text_parts.extend(get_shape_text(sub_shape))
    return text_parts

def get_shape_images(shape):
    images = []
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        if hasattr(shape, "image"):
            images.append(shape.image.blob)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            images.extend(get_shape_images(sub_shape))
    return images

def generate_questions_api(content, slide_id, api_key):
    if not content.strip():
        return "No content to generate questions."
    
    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
    prompt = (
        f"Analyze the content (ID: {slide_id}). Generate a list of potential user search queries "
        "that this content answers. Focus on 'What is', 'How to', 'Explain'. "
        "Return ONLY the questions list."
        f"\n\nContent:\n{content}"
    )
    payload = {
        "model": "deepseek-chat",
        "messages": [
            {"role": "system", "content": "You are a helpful AI."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7
    }
    try:
        resp = requests.post(DEEPSEEK_URL, headers=headers, json=payload)
        resp.raise_for_status()
        return resp.json()['choices'][0]['message']['content'].strip()
    except Exception as e:
        return f"API Error: {e}"

def process_ppt_file(uploaded_file, tesseract_path, api_key):
    # Setup Tesseract
    pytesseract.pytesseract.tesseract_cmd = tesseract_path
    
    # Save uploaded file to temp (Use Absolute Path for COM compatibility)
    abs_path = os.path.abspath("temp_input.pptx")
    with open(abs_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    st.session_state.current_ppt_path = abs_path
    
    ppt_filename = uploaded_file.name
    prs = Presentation(abs_path)
    output_data = []
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_id = f"{ppt_filename}_slide_{slide_num}"
        status_text.text(f"Processing Slide {slide_num}/{total_slides}...")
        
        # Extract Text & Images
        extracted_parts = []
        
        for shape in slide.shapes:
            extracted_parts.extend(get_shape_text(shape))
            for blob in get_shape_images(shape):
                try:
                    img = Image.open(io.BytesIO(blob))
                    text = pytesseract.image_to_string(img, config='--psm 6')
                    if text.strip():
                        extracted_parts.append(f"[Image Text]: {text.strip()}")
                except: pass
        
        full_text = "\n".join(extracted_parts).strip()
        
        # Generate Questions
        questions = generate_questions_api(full_text, slide_id, api_key)
        
        # Dynamic Keys Logic
        text_key = f"text_{slide_id}"
        q_key = f"Questions_{slide_id}"

        output_data.append({
            "slide": slide_num,
            "ppt_id": ppt_filename,
            "slideID": slide_id,
            text_key: full_text,
            q_key: questions
        })
        
        progress_bar.progress((i + 1) / total_slides)

    status_text.text("Processing Complete!")
    return output_data

# ==========================================
# 4. VECTOR DB LOGIC
# ==========================================

def build_vector_db(data):
    client = chromadb.PersistentClient(path="./chroma_db")
    
    try:
        client.delete_collection("ppt_data")
    except ValueError:
        pass 

    model = SentenceTransformer(EMBEDDING_MODEL_NAME)
    
    class LocalEmbeddingFunction(chromadb.EmbeddingFunction):
        def __call__(self, input):
            return model.encode(input).tolist()

    embedding_fn = LocalEmbeddingFunction()
    collection = client.create_collection(name="ppt_data", embedding_function=embedding_fn)

    documents = []
    metadatas = []
    ids = []

    for item in data:
        s_id = item['slideID']
        text_content = item.get(f"text_{s_id}", "")
        questions_content = item.get(f"Questions_{s_id}", "")
        
        documents.append(questions_content)
        
        metadatas.append({
            "slideID": s_id, 
            "ppt_id": item['ppt_id'], 
            "page_num": item['slide'],
            "slide_text": text_content
        })
        ids.append(s_id)

    collection.add(documents=documents, metadatas=metadatas, ids=ids)
    return client, collection

def query_vector_db(query, collection, n_results=1):
    results = collection.query(query_texts=[query], n_results=n_results)
    return results

def generate_slide_answer(query, slide_text, api_key):
    prompt = (
        "You are a strict analyst. Your task is to answer the user's question based ONLY on the provided slide content.\n\n"
        "Rules:\n"
        "1. Check if the provided Slide Content contains the information to answer the User Question.\n"
        "2. If the content is relevant and contains the answer, output the answer based strictly on the text.\n"
        "3. If the content is NOT relevant or does not contain the answer, output exactly the word 'IRRELEVANT'.\n"
        "4. Do NOT use outside knowledge. Do NOT hallucinate.\n\n"
        f"Slide Content:\n{slide_text}\n\n"
        f"User Question: {query}"
    )
    
    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
    payload = {
        "model": "deepseek-chat",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.1 
    }
    
    try:
        resp = requests.post(DEEPSEEK_URL, headers=headers, json=payload)
        return resp.json()['choices'][0]['message']['content'].strip()
    except Exception as e:
        return f"Error: {e}"

# ==========================================
# 5. STREAMLIT UI LAYOUT
# ==========================================

# --- SIDEBAR (CONFIG + VISUAL SLIDE VIEWER) ---
with st.sidebar:
    # 1. SIDEBAR VIEWER LOGIC (Top Priority)
    if st.session_state.active_slide_id:
        with st.container(border=True):
            st.subheader(f"📑 Viewer: {st.session_state.active_slide_id}")
            
            # Close Button
            if st.button("❌ Close Viewer", use_container_width=True):
                st.session_state.active_slide_id = None
                st.rerun()

            # Find Slide Index
            slide_idx = 1
            if st.session_state.extracted_data:
                for item in st.session_state.extracted_data:
                    if item["slideID"] == st.session_state.active_slide_id:
                        slide_idx = item["slide"]
                        break
            
            # --- RENDER SLIDE VISUALLY ---
            if st.session_state.current_ppt_path and os.path.exists(st.session_state.current_ppt_path):
                with st.spinner("Rendering Slide View..."):
                    img_path = render_slide_as_image(st.session_state.current_ppt_path, slide_idx)
                    
                    if img_path and os.path.exists(img_path):
                        st.image(img_path, caption=f"Original Slide {slide_idx}", use_container_width=True)
                    else:
                        st.error("Could not render slide visual. (PowerPoint may not be accessible).")
                        # Fallback to text
                        st.warning("Showing extracted text instead:")
                        slide_text = ""
                        for item in st.session_state.extracted_data:
                            if item["slideID"] == st.session_state.active_slide_id:
                                slide_text = item.get(f"text_{st.session_state.active_slide_id}", "")
                                break
                        st.text_area("Content", slide_text, height=200, disabled=True)
            else:
                st.error("PPT Source file missing.")
                
        st.divider()

    # 2. CONFIGURATION
    st.header("⚙️ Configuration")
    api_key_input = st.text_input("DeepSeek API Key", value=DEFAULT_API_KEY, type="password")
    tesseract_input = st.text_input("Tesseract Path", value=DEFAULT_TESSERACT_PATH)
    
    st.divider()
    
    uploaded_file = st.file_uploader("Upload PowerPoint", type=["pptx"])
    
    if uploaded_file and st.button("🚀 Process & Ingest"):
        with st.spinner("Extracting Text & Images... (This may take a moment)"):
            data = process_ppt_file(uploaded_file, tesseract_input, api_key_input)
            st.session_state.extracted_data = data
            
        with st.spinner("Building Vector Database..."):
            build_vector_db(data)
            st.session_state.vector_db_ready = True
            st.success("Ingestion Complete!")

# --- MAIN PAGE ---
st.title("📄 Intelligent PPT Pipeline Explorer")

tab1, tab2, tab3 = st.tabs(["🔍 Search Pipeline", "💬 Chat Assistant", "📊 Extracted Data"])

# TAB 1: PIPELINE INTERFACE
with tab1:
    if not st.session_state.vector_db_ready:
        st.info("👈 Please upload a PPT and click 'Process & Ingest' to start.")
    else:
        client = chromadb.PersistentClient(path="./chroma_db")
        model = SentenceTransformer(EMBEDDING_MODEL_NAME)
        class LocalEmbeddingFunction(chromadb.EmbeddingFunction):
            def __call__(self, input):
                return model.encode(input).tolist()
        collection = client.get_collection("ppt_data", embedding_function=LocalEmbeddingFunction())

        st.subheader("Run Query Pipeline")
        top_k = st.slider("Select Top-K Results", min_value=1, max_value=5, value=1)
        user_query = st.chat_input("Enter your question here...", key="pipeline_input")
        
        if user_query:
            st.markdown(f"**Question:** {user_query}")
            st.divider()
            st.info(f"### Phase 1: Vector Search (Retrieving Top {top_k} Matches)")
            results = query_vector_db(user_query, collection, n_results=top_k)
            
            num_matches = len(results['ids'][0])
            for i in range(num_matches):
                matched_slide_id = results['ids'][0][i]
                matched_questions = results['documents'][0][i]
                similarity_distance = results['distances'][0][i]
                retrieved_text = results['metadatas'][0][i]['slide_text']
                
                st.markdown(f"#### Result #{i+1} (Slide ID: `{matched_slide_id}`)")
                col1, col2 = st.columns([3, 1])
                with col1:
                    st.write("**Matched Synthetic Questions:**")
                    st.code(matched_questions, language="text")
                with col2:
                    st.metric("Match Distance", f"{similarity_distance:.4f}")
                st.warning(f"**Phase 2: Context Retrieval (text_{matched_slide_id})**")
                st.text_area(f"Retrieved Slide Text (#{i+1}):", retrieved_text, height=150)
                st.divider()

# TAB 2: CHAT ASSISTANT
with tab2:
    st.subheader("🤖 Slide-by-Slide Assistant")
    
    if not st.session_state.vector_db_ready:
        st.info("👈 Please process a document first.")
    else:
        chat_k = st.slider("Context Slides to Check", 1, 5, 3, key="chat_k_slider")
        
        # Display Chat History with Interactivity
        for idx, msg in enumerate(st.session_state.chat_history):
            with st.chat_message(msg["role"]):
                if "structured_results" in msg:
                    st.markdown(msg["content"])
                    st.caption("🔗 Relevant Slides (Click to View in Sidebar):")
                    cols = st.columns(len(msg["structured_results"]))
                    for col_i, item in enumerate(msg["structured_results"]):
                        s_id = item['slide_id']
                        if cols[col_i].button(f"📄 View {s_id}", key=f"hist_btn_{idx}_{s_id}"):
                            st.session_state.active_slide_id = s_id
                            st.rerun()
                else:
                    st.markdown(msg["content"])
                
        if prompt := st.chat_input("Ask a question about your presentation...", key="chat_input"):
            st.session_state.chat_history.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)
            
            with st.chat_message("assistant"):
                message_placeholder = st.empty()
                message_placeholder.markdown("Searching & Analyzing...")
                
                client = chromadb.PersistentClient(path="./chroma_db")
                model = SentenceTransformer(EMBEDDING_MODEL_NAME)
                class LocalEmbeddingFunction(chromadb.EmbeddingFunction):
                    def __call__(self, input):
                        return model.encode(input).tolist()
                collection = client.get_collection("ppt_data", embedding_function=LocalEmbeddingFunction())
                
                results = query_vector_db(prompt, collection, n_results=chat_k)
                num_matches = len(results['ids'][0])
                
                structured_results = []
                full_response_text = ""
                
                for i in range(num_matches):
                    s_id = results['ids'][0][i]
                    s_text = results['metadatas'][0][i]['slide_text']
                    
                    answer = generate_slide_answer(prompt, s_text, api_key_input)
                    
                    if answer.upper() != "IRRELEVANT" and "IRRELEVANT" not in answer:
                        entry = f"**{s_id}**\n{answer}\n\n---\n\n"
                        full_response_text += entry
                        structured_results.append({'slide_id': s_id, 'answer': answer})
                
                if not structured_results:
                    full_response_text = "I couldn't find any relevant information in the top search results."
                
                message_placeholder.markdown(full_response_text)
                
                if structured_results:
                    st.caption("🔗 Relevant Slides (Click to View in Sidebar):")
                    cols = st.columns(len(structured_results))
                    for col_i, item in enumerate(structured_results):
                        s_id = item['slide_id']
                        if cols[col_i].button(f"📄 View {s_id}", key=f"new_btn_{s_id}"):
                            st.session_state.active_slide_id = s_id
                            st.rerun()
                
                st.session_state.chat_history.append({
                    "role": "assistant", 
                    "content": full_response_text,
                    "structured_results": structured_results
                })

# TAB 3: DATA INSPECTION
with tab3:
    if st.session_state.extracted_data:
        st.write(f"Total Slides: {len(st.session_state.extracted_data)}")
        st.json(st.session_state.extracted_data)
    else:
        st.write("No data processed yet.")
